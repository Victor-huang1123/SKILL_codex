"""
Generate a PowerPoint deck using the January group-meeting PPTX as a real template.

Key idea:
- Use slide layouts/placeholders from group_meeting_20260108.pptx.
- Put title text into the title placeholder.
- Put bullet text into the content placeholder.
- Use paragraph.level for bullet hierarchy.
- Do NOT type bullet symbols manually, e.g., no "• xxx".

Required package:
    pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path
from typing import Literal, Sequence

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# -----------------------------
# Skill paths
# -----------------------------
SKILL_DIR = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = SKILL_DIR / "assets" / "template_python.pptx"
OUTPUT_PATH = Path("output_template_driven.pptx")

# Layout indices observed in group_meeting_20260108.pptx:
# 0: 標題投影片
# 1: 標題及物件      <- title + content placeholder; best for normal bullet slides
# 3: 兩項物件        <- title + two content placeholders
# 5: 只有標題        <- title only; useful for custom figure slides
# 6: 空白            <- blank; use only when absolutely needed
TITLE_LAYOUT = 0
CONTENT_LAYOUT = 1
TWO_CONTENT_LAYOUT = 3
TITLE_ONLY_LAYOUT = 5

FooterMode = Literal["none", "template", "manual"]
BulletKind = Literal["bullet", "plain", "equation"]
BulletItem = tuple[int, str] | tuple[int, str, BulletKind]
FigureBox = tuple[int, int, int, int]

DEFAULT_LEFT_TEXT_BOX = (
    Inches(0.72),
    Inches(1.34),
    Inches(5.25),
    Inches(4.85),
)
DEFAULT_RIGHT_FIGURE_BOX = (
    Inches(6.15),
    Inches(1.33),
    Inches(4.98),
    Inches(4.45),
)


# -----------------------------
# Low-level helpers
# -----------------------------
def delete_all_slides(prs: Presentation) -> None:
    """Remove all existing slides while keeping slide masters/layouts.

    python-pptx has no public API for deleting slides, so this uses the standard
    internal slide-id list method. This is common when using a PPTX as a template.
    """
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    rel_ids = [slide_id.rId for slide_id in slide_id_list]
    for rel_id in rel_ids:
        prs.part.drop_rel(rel_id)
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)


def placeholder(slide, idx: int):
    """Return a placeholder by idx, or None if it does not exist."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def set_footer(slide, slide_no: int, mode: FooterMode = "manual") -> None:
    """Fill January-template footer placeholders.

    idx 10 is the GIEE.NTU/date footer in the January template.
    idx 11 or 12 is the slide number placeholder depending on layout.
    """
    if mode == "none":
        return

    if mode == "manual":
        footer = placeholder(slide, 10)
        if footer is not None and hasattr(footer, "text"):
            footer.text = " GIEE.NTU"

        number = placeholder(slide, 11) or placeholder(slide, 12)
        if number is not None and hasattr(number, "text"):
            number.text = str(slide_no)


def remove_bullet(paragraph) -> None:
    """Force a paragraph to have no bullet via OOXML.

    Useful for equation lines inside a content placeholder.
    """
    p_pr = paragraph._p.get_or_add_pPr()  # noqa: SLF001

    # Remove existing bullet settings if present.
    for child in list(p_pr):
        if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
            p_pr.remove(child)

    # Add <a:buNone/>.
    from pptx.oxml.xmlchemy import OxmlElement

    p_pr.append(OxmlElement("a:buNone"))


def style_equation(paragraph) -> None:
    """Apply lightweight direct formatting only for formula-like lines."""
    remove_bullet(paragraph)
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    for run in paragraph.runs:
        run.font.name = "Cambria Math"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def normalize_text_frame(text_frame) -> None:
    """Keep text inside its placeholder without PowerPoint auto-shrinking it."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE


def set_shape_box(shape, box: FigureBox) -> None:
    """Move and resize a placeholder or shape to an EMU box."""
    left, top, width, height = box
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height


def add_picture_fit(slide, image_path: str | Path, box: FigureBox = DEFAULT_RIGHT_FIGURE_BOX):
    """Insert an image and fit it inside a box without changing aspect ratio."""
    left, top, width, height = box
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(width / pic.width, height / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    return pic


def add_caption(slide, caption: str, box: FigureBox) -> None:
    """Add a compact caption under a figure."""
    if not caption:
        return
    textbox = slide.shapes.add_textbox(*box)
    tf = textbox.text_frame
    tf.clear()
    normalize_text_frame(tf)
    p = tf.paragraphs[0]
    p.text = caption
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


# -----------------------------
# Template-driven slide builders
# -----------------------------
def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    footer_mode: FooterMode = "manual",
):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    slide_no = len(prs.slides)

    title_ph = placeholder(slide, 0) or slide.shapes.title
    if title_ph is not None:
        title_ph.text = title

    subtitle_ph = placeholder(slide, 1)
    if subtitle_ph is not None and hasattr(subtitle_ph, "text"):
        subtitle_ph.text = subtitle

    set_footer(slide, slide_no, footer_mode)
    return slide


def add_bullet_slide(
    prs: Presentation,
    title: str,
    items: Sequence[BulletItem],
    footer_mode: FooterMode = "manual",
):
    """Create a normal title + content slide using January template layout 1.

    items format:
        (level, text)                  -> native bullet, level 0/1/2...
        (level, text, "bullet")        -> native bullet
        (level, text, "plain")         -> no bullet, inherited font
        (level, text, "equation")      -> no bullet, Cambria Math 18 pt
    """
    slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
    slide_no = len(prs.slides)

    title_ph = placeholder(slide, 0) or slide.shapes.title
    if title_ph is not None:
        title_ph.text = title

    body = placeholder(slide, 1)
    if body is None or not hasattr(body, "text_frame"):
        raise RuntimeError("Content placeholder idx=1 not found. Check template layout 1.")

    tf = body.text_frame
    tf.clear()
    normalize_text_frame(tf)

    for i, raw_item in enumerate(items):
        if len(raw_item) == 2:
            level, text = raw_item  # type: ignore[misc]
            kind: BulletKind = "bullet"
        else:
            level, text, kind = raw_item  # type: ignore[misc]

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level

        if kind == "plain":
            remove_bullet(p)
        elif kind == "equation":
            style_equation(p)
        # For normal bullets: do NOT set font name/size/color here.
        # Let the January PowerPoint template/master control bullet formatting.

    set_footer(slide, slide_no, footer_mode)
    return slide


def add_outline_slide(
    prs: Presentation,
    sections: Sequence[str],
    current: str | None = None,
    footer_mode: FooterMode = "manual",
) -> object:
    """Add an Outline slide with template-native bullets.

    If current is given, it is still kept as normal text to avoid hard-coded styling.
    You can manually emphasize it later in PowerPoint if needed.
    """
    items: list[BulletItem] = []
    for section in sections:
        items.append((0, section))
    return add_bullet_slide(prs, "Outline", items, footer_mode=footer_mode)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_items: Sequence[BulletItem],
    right_items: Sequence[BulletItem],
    footer_mode: FooterMode = "manual",
):
    """Use January template layout 3: title + two content placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT_LAYOUT])
    slide_no = len(prs.slides)

    title_ph = placeholder(slide, 0) or slide.shapes.title
    if title_ph is not None:
        title_ph.text = title

    for ph_idx, items in [(1, left_items), (2, right_items)]:
        body = placeholder(slide, ph_idx)
        if body is None or not hasattr(body, "text_frame"):
            continue
        tf = body.text_frame
        tf.clear()
        normalize_text_frame(tf)
        for i, raw_item in enumerate(items):
            if len(raw_item) == 2:
                level, text = raw_item  # type: ignore[misc]
                kind: BulletKind = "bullet"
            else:
                level, text, kind = raw_item  # type: ignore[misc]
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = level
            if kind == "plain":
                remove_bullet(p)
            elif kind == "equation":
                style_equation(p)

    set_footer(slide, slide_no, footer_mode)
    return slide


def add_figure_slide(
    prs: Presentation,
    title: str,
    items: Sequence[BulletItem],
    image_path: str | Path,
    caption: str = "",
    text_box: FigureBox = DEFAULT_LEFT_TEXT_BOX,
    figure_box: FigureBox = DEFAULT_RIGHT_FIGURE_BOX,
    footer_mode: FooterMode = "manual",
):
    """Create a bullet slide with a fitted figure/table on the right."""
    slide = add_bullet_slide(prs, title, items, footer_mode=footer_mode)
    body = placeholder(slide, 1)
    if body is not None:
        set_shape_box(body, text_box)
    add_picture_fit(slide, image_path, figure_box)
    caption_box = (
        figure_box[0],
        figure_box[1] + figure_box[3] + Inches(0.06),
        figure_box[2],
        Inches(0.32),
    )
    add_caption(slide, caption, caption_box)
    return slide


# -----------------------------
# Example deck content
# -----------------------------
def build_example_deck(template_path: Path = TEMPLATE_PATH, output_path: Path = OUTPUT_PATH) -> None:
    prs = Presentation(str(template_path))
    delete_all_slides(prs)

    sections = [
        "Introduction",
        "Preliminaries",
        "Problem Formulation",
        "Methodology",
        "Experimental Results",
        "Conclusion",
    ]

    add_title_slide(
        prs,
        title="A Graph-Based Approach for Optimizing Pin Access in Nanosheet FET Standard Cell Library Synthesis",
        subtitle=(
            "2026-05-27\n"
            "Presenter: Zhih-Chi Huang\n"
            "The Electronic Design Automation Laboratory\n"
            "Graduate Institute of Electronics Engineering\n"
            "National Taiwan University"
        ),
    )

    add_outline_slide(prs, sections)

    add_bullet_slide(
        prs,
        "Problem Formulation",
        [
            (0, "Given"),
            (1, "Gate-level netlist with PMOS/NMOS transistor information"),
            (1, "Nanosheet design rules and routing constraints"),
            (1, "Candidate M0/M1 pin-access locations"),
            (0, "Output"),
            (1, "Legal single-row or multi-row standard cell layout"),
            (1, "Routed internal connections"),
            (1, "I/O pins assigned on M0 or M1"),
            (0, "Objective"),
            (1, "Minimize placement and routing cost"),
            (1, "Maximize accessible pin paths"),
            (1, "Avoid pin overlap and excessive pin extension"),
        ],
    )

    add_bullet_slide(
        prs,
        "Graph-Based Routing Model",
        [
            (0, "Each two-pin subnet is routed on a connection graph"),
            (1, "Multi-pin nets are decomposed into two-pin subnets"),
            (1, "Pin vertices are fixed to degree 1 by routing constraints"),
            (0, "A physical graph keeps different signals mutually exclusive"),
            (1, "Physical graph models all cell-layout routing resources"),
            (1, "Mutual-exclusion constraints prevent illegal sharing across signals"),
            (0, "Steiner-tree reuse allows the same net to share routing resources"),
            (1, "Non-pin vertices are restricted to degree 0 or 2"),
            (1, "Resources can be reused only within the same multi-pin net"),
        ],
    )

    add_bullet_slide(
        prs,
        "Conclusion",
        [
            (0, "First standard-cell synthesis flow tailored to nanosheet FET with M0 / M1 pin allocation"),
            (0, "Virtual node modeling estimates boundary access resources"),
            (0, "Pin overlap reduction and dynamic pin metal selection improve access quality"),
            (0, "Transistor coarsening keeps the SMT model tractable"),
            (0, "Block-level P&R improves area, DRV, and wirelength simultaneously"),
            (1, "vs NS3K at 70% utilization: -3.2% area, -97.6% DRV, -16.5% wirelength"),
            (1, "vs LP heuristic: -77.1% DRV at no wirelength cost"),
        ],
    )

    prs.save(str(output_path))
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    build_example_deck()
